from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE

import re

from model import db, connect_to_db, Clip, TextPull

from moviepy.editor import *

# set horizontal and vertical position
vid_VP_left = Inches(1.91)
vid_HP_left = Inches(.41)

vid_HP_center = Inches(2.5)
vid_VP_center = Inches(1.22)

# set video width and height
vid_H_sm = Inches(2.82)
vid_W_sm = Inches(3.77)

vid_H_lg = Pt(400)
vid_W_lg = Pt(600)

def create_slide_deck(template, clips, vid_name, curr_time):
    """Creates a slide deck of selected clips"""

    # create a new presentation
    prs = Presentation(template)

    pres_name = vid_name + "_" + curr_time + ".pptx"

    # make a new slide for each clip
    for clip in clips:
        # pull out just the video name without the file path
        clip_name = clip[clip.rfind('/')+1:]

        # get the clip from the db
        db_clip = db.session.query(Clip).filter(Clip.clip_name == clip_name).first()

        # get a still from that video
        vid_clip = VideoFileClip(clip)

        # make a png of the first frame of the video to be the front
        vid_front = clip[:-4]+".png"
        vid_png = vid_clip.save_frame(vid_front)

        deponent = db_clip.video.deponent

        if db_clip.textpull:
            # select the depo slide layout
            slide = prs.slides.add_slide(prs.slide_layouts[7])

            # add a video to the slide using the provided dimensions
            video = slide.shapes.add_movie(clip, vid_HP_left, vid_VP_left, vid_W_sm, vid_H_sm, vid_front)

            # get the matching textpull for that clip
            db_text = db.session.query(TextPull).filter(TextPull.clip_id == db_clip.clip_id).first()

            # we have to split the lines up into pieces so we can insert them properly
            # if it is a solid string - we lose our hard returns
            split_text = db_text.pull_text.split("\n")

            citation = u'\u2013' + 'Depo at ' + db_clip.start_pl + u'\u2013' + db_clip.end_pl

            # target the text frame on the slide
            txt = slide.placeholders[1]

            cite = slide.placeholders[15]
            cite.text = citation

            if txt.has_text_frame:
                # find the text frame in the placeholder
                text_frame = txt.text_frame

                # set text frame to be auto-sizing
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                # id the paragraphs
                p = text_frame.paragraphs[0]

                # entries starting with a Q. or A. have a blank space at the start
                # so if the 0th "line" is nothing - start at line 1
                start_line = 0
                if split_text[0] != "":
                    p.text = split_text[0]
                else:
                    p.text = split_text[1]
                    start_line = 1

                # load the text in line by line as paragraphs
                for para in split_text[start_line+1:]:
                    p = text_frame.add_paragraph()
                    p.text = para
        else:
            # select the 1-line blank slide master
            slide = prs.slides.add_slide(prs.slide_layouts[10])

            # add a video to the slide using the provided dimensions
            video = slide.shapes.add_movie(clip, vid_HP_center, vid_VP_center, vid_W_lg, vid_H_lg, vid_front)

        # set deponent name as title
        title_placeholder = slide.shapes.title
        title_placeholder.text = deponent

    prs.save('static/ppts/'+pres_name)


def find_text_by_page_line(start_pl, end_pl, source):
    """Pulls text based on page and line numbers"""

    # split source text into a list - it goes in as readlines()
    source = source.split('","')

    # split the page num from the line num
    start_page_line = start_pl.split(':')
    end_page_line = end_pl.split(':')

    # convert to int
    start_page = int(start_page_line[0])
    start_line = int(start_page_line[1])

    # if the end page is the same as the start page
    # set end page to be start page
    if len(end_page_line) == 1:
        end_page = start_page
        end_line = int(end_page_line[0])
    else:
        end_page = int(end_page_line[0])
        end_line = int(end_page_line[1])

    line_total = 0

    # if the lines are on the same page
    if start_page == end_page:
        line_total = end_line - start_line + 1
    # if pull spans more than 2 pages
    elif start_page + 1 != end_page:
        # find the number of full pages btwn start and end pages
        middle_pages = end_page - start_page - 1
        line_total = (26-start_line) + (middle_pages * 25) + end_line
    # if they aren't - collect the line total from first page and end line from last
    else:
        line_total = (26-start_line) + end_line

    # make a var to hold starting position of text
    start_pos = 0

    # search the lines of the transcript for the page number
    for i in range(len(source)):
        result = re.search('00%s:\d{2}' % start_page, source[i])
        # once you find it - add the staring line to the i index
        # so if it was page 13, line 5 - find 013:01
        # then add 4 (5-1) lines to it - we subtract 1 bc we start on line 1
        if result is not None:
            start_pos = i + start_line - 1
            break

    end_pos = line_total+start_pos

    print start_pos, end_pos
    # make a list to hold the lines of text
    cropped_text = [];
    start_at = ""
    end_at = ""

    timecode_only = re.compile("\s{2}(\d{2}:\d{2}:\d{2}.\d{3})\s")
    # starting from the start_position and spanning the number of requested lines
    # add each selected line to the list for clean up
    for i in range(start_pos, end_pos):
        if i == start_pos:
            start_at = re.search(timecode_only, source[i]).group(1)
        if i == end_pos - 1:
            end_at = re.search(timecode_only, source[i]).group(1)
        cropped_text.append(source[i])

    # remove page and line excess 
    pulled_text = clean_up_transcript_text(cropped_text)

    # fix punctuation
    pulled_text = fix_punctuation(pulled_text)

    return (start_at, end_at, pulled_text)

def clean_up_transcript_text(transcript_lines):
    """Removes page/line/timecodes from raw transcript, returns clean string"""

    # ex:   00:15:59.946        17        or  00:16:46.550  00019:01       
    page_line_pattern = re.compile("\s{2}\d{2}:\d{2}:\d{2}.\d{3}\s*\d{2}\s*|\d{3}:\d{2}\s*")

    # find a capital A or Q followed by 5 spaces
    q_a_pattern = re.compile("([QA])\s{5}")

    clean_text = []

    for line in transcript_lines:
        line = line.rstrip()
        clean_line = re.sub(page_line_pattern, "", line)
        # print clean_line
        clean_line = re.sub(q_a_pattern, '\g<1>.\t', clean_line)
        # if its a Q. or A. line - add a new line before it
        if clean_line.startswith("Q.") or clean_line.startswith("A."):
            clean_line = "\n" + clean_line
        # add to list of cleaned up text
        clean_text.append(clean_line)

    # print clean_text
    pulled_text = " ".join(clean_text)

    return pulled_text


def fix_punctuation(pulled_text):
    """Fixes dumb quotes, em dashes, extra spaces"""

    # finds quote marks with a preceeding space which would be open
    open_quote = re.compile('\s"')

    # finds em dashes with spaces on either side or both sides
    em_dash = re.compile(' ?-- ?')

    # fix open quotes, and then any remaining dumb double quotes to be closing
    pulled_text = re.sub(open_quote, u'\u201C', pulled_text)
    pulled_text = re.sub('"', u'\u201D', pulled_text)

    # fix em dashes
    pulled_text = re.sub(em_dash, u'\u2014', pulled_text)

    # fix single quotes
    pulled_text = re.sub("'", u'\u2019', pulled_text)

    # finally remove double spaces
    pulled_text = re.sub("  ", " ", pulled_text)

    return pulled_text
